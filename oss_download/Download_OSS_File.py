import asyncio
import aiohttp
import pymysql
import os
import logging
from tenacity import retry, stop_after_attempt, wait_fixed
from openpyxl import Workbook, load_workbook
import m3u8
import subprocess
from pathlib import Path
from tqdm import tqdm  # 用于进度条显示
import time
from urllib.parse import urlparse
import alibabacloud_oss_v2 as oss
from alibabacloud_oss_v2.models import PutObjectRequest
from datetime import datetime
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image

###############################################################
# 用来爬取oss库的内容的，因为分不开，没法oss对oss的迁移，只能爬取后上传
# 
# 2个已知不完美的地方
# 1、针对mp4格式的文件，可能存在丢帧情况，视频没损坏，但是下载后可能比如30分钟，中间有播放直接到结尾了，尾部有部分数据没下载完整，可能跟分块下载有关系，或者分块下载修改成1024 KB，不要用成1M，待测试和验证，发生几率不高
# 2、获取ppt和pptx文档的页数的时候，如果设置了页面隐藏，可能页数获取不准，有代码解决，但是未测试和验证
# 
# 其他问题：
# 1、并发爬取效率没觉得快多少
# 2、中间有段根据后缀判断的逻辑弯弯绕绕的是因为规则很不规律，一会带文件名后缀，一会不带文件名后缀，麻烦，所以写了很多if esle强制判断
# 
# 注意：
# 1、需要安装ffmpeg  流媒体处理方式，视频文件也可以处理，但是速度比较慢
# 2、需要安装aria2   比分块下载稳定性好
#
# 整体爬虫思路可以供参考
###############################################################

# 去掉ppt的隐藏页代码如下，未验证 供参考
def count_visible_slides(pptx_file):
    presentation = Presentation(pptx_file)
    visible_slide_count = 0

    for slide in presentation.slides:
        # 检查幻灯片的可见性
        if not slide.hidden:  # 检查幻灯片是否隐藏
            visible_slide_count += 1

    return visible_slide_count

# 1、分块下载可能有几率丢帧的问题也可以用ffmpeg解决，就是速度比较慢
# yum install epel-release
# yum install ffmpeg ffmpeg-devel
# ffmpeg -i "http://example.com/path/to/video.mp4" -c copy output.mp4

# 2、如果需要异步下载可以求助于aria2
# 求助于安装独立下aria2下载工具提高速度和减少丢帧可靠性
# yum install aria2
def download_with_aria2(url, filename):
    subprocess.run(['aria2c', url, '-o', filename])

url = 'https://example.com/path/to/video.mp4'  # 替换为视频的实际 URL
download_with_aria2(url, 'video.mp4')
     


# 配置字典
config = {
    "db": {
        "host": "127.0.0.1",
        "user": "username",
        "password": "password",
        "database": "dbname",
        "port": 3306
    },
    "oss": {
        "enable_upload": False,                          # False,  # 是否启用上传到 OSS
        "access_key_id": "id",
        "access_key_secret": "secret",
        "bucket_name": "bucket",
        "region": "cn-hangzhou",                         # Bucket所在区域
        "endpoint": f"oss-{region}.aliyuncs.com",
        "base_url": f"https://{bucket_name}.{endpoint}"  # OSS 访问的基础 URL
    },
    "download": {
        "save_path": "./downloads",  # 本地保存路径
        "concurrent_limit": 100,     # 最大并发下载数
        "retry_attempts": 3,         # 最大重试次数
        "retry_interval": 5,         # 重试间隔时间（秒）
        "chunk_size": 1024,          # 文件分块大小（1024*1024,1MB  有几率存在丢帧情况）
        "time_out": 3600 * 8         # 请求超时时间（秒）因为有几个G的
    },
    "excel": {
        "enable_update_db": False,          # 是否修改数据库
        "progress_file": "./progress.xlsx"  # 下载进度文件路径
    },
    "m3u8": {
        "ffmpeg_path": "ffmpeg",    # ffmpeg 可执行文件路径
        "temp_path": "./temp_m3u8"  # m3u8 临时文件存储路径
    }
}

# 查询的SQL
SQL_QUERY="""
select id, resource_id, resource_name, resource_path, resource_type_tag, create_time from dim_resource_info 
where is_downloaded <> 1 and resource_path like 'http%'
"""


# 初始化日志
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler("download.log", encoding="utf-8"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# 数据库连接
def get_db_connection():
    return pymysql.connect(
        host=config["db"]["host"],
        user=config["db"]["user"],
        password=config["db"]["password"],
        database=config["db"]["database"],
        port=config["db"]["port"],
        cursorclass=pymysql.cursors.DictCursor
    )

# 从数据库读取素材数据
def fetch_material_data():
    conn = get_db_connection()
    try:
        with conn.cursor() as cursor:
            cursor.execute(SQL_QUERY)
            return cursor.fetchall()
    finally:
        conn.close()

# 更新数据库记录
def update_resource_status(resource_id, is_downloaded, download_time=None, is_upload_oss=None, upload_time=None, oss_path=None, error_message=None):
    if not config["excel"]["enable_update_db"]:
        logger.info(f"update DB status 已禁用，跳过：{resource_id}")
        return

    conn = get_db_connection()
    try:
        with conn.cursor() as cursor:
            sql = "UPDATE dim_resource_info SET is_downloaded=%s, download_time=%s, is_upload_oss=%s, upload_time=%s, oss_path=%s, error_message=%s WHERE resource_id=%s"
            cursor.execute(sql, (is_downloaded, download_time, is_upload_oss, upload_time, oss_path, error_message, resource_id))
            conn.commit()
    except Exception as e:
        logger.error(f"更新数据库状态失败：{e}")
    finally:
        conn.close()

# 更新数据库记录文档页数
def update_resource_page_count(resource_id, page_count):
    if not config["excel"]["enable_update_db"]:
        logger.info(f"update DB pagecount 已禁用，跳过：{resource_id}")
        return

    conn = get_db_connection()
    try:
        with conn.cursor() as cursor:
            sql = "UPDATE dim_resource_info SET page_count=%s WHERE resource_id=%s"
            cursor.execute(sql, (page_count, resource_id))
            conn.commit()
    except Exception as e:
        logger.error(f"更新数据库文档页数失败：{e}")
    finally:
        conn.close()

# 校验 resourcepath 的合法性
def is_valid_resourcepath(resourcepath):
    """
    检查 URL 是否是合法的 HTTP/HTTPS 地址
    :param url: 要检查的 URL
    :return: 如果合法返回 True，否则返回 False
    """
    try:
        parsed = urlparse(resourcepath)
        return parsed.scheme in ('http', 'https') and parsed.netloc != ''
    except Exception as e:
        logger.warning(f"URL 无效: {resourcepath}, 错误信息: {e}")
        return False

# 初始化 OSS 客户端
def get_alicloud_oss_client():
    # 将字典中的凭证写入当前进程的环境变量
    os.environ["OSS_ACCESS_KEY_ID"] = config['oss']['access_key_id']
    os.environ["OSS_ACCESS_KEY_SECRET"] = config['oss']['access_key_secret']

    # 初始化客户端（自动从环境变量加载）
    credentials_provider = oss.credentials.EnvironmentVariableCredentialsProvider()
    cfg = oss.config.load_default()
    cfg.credentials_provider = credentials_provider

    cfg.region = config['oss']['region']  # 设置区域[citation:1]
    cfg.endpoint = config['oss']['endpoint']

    # 初始化OSS客户端[citation:1]
    return oss.Client(cfg)

# 常规文件上传
def upload_simple_file_to_oss(local_file_path, object_name=None):
    """
    使用阿里云OSS Python SDK V2上传文件到OSS

    参数:
        local_file_path: 本地文件路径
        object_name: OSS中保存的文件路径，如果为None则使用本地文件名
    """
    # 如果未指定OSS对象名称，则使用本地文件名
    if object_name is None:
        object_name = os.path.basename(local_file_path)

    try:
        # 获取aliyun oss client
        client = get_alicloud_oss_client()

        # 使用简单上传接口上传文件[citation:3][citation:4]
        with open(local_file_path, 'rb') as file:
            result = client.put_object_from_file(
                PutObjectRequest(
                    bucket=config['oss']['bucket_name'],
                    key=object_name,
                    body=file
                ),
                local_file_path          # 本地文件路径
            )

        # 构建完整的访问URL
        full_url = f"{config['oss']['base_url']}/{object_name}"
        logger.info(f"文件上传成功:{result.status_code} - 文件访问URL: {full_url}")

        return {
            'success': True,
            'etag': result.etag,
            'url': full_url,
            'object_name': object_name
        }

    except Exception as e:
        logger.error(f"文件上传失败: {str(e)}")
        return {
            'success': False,
            'error': str(e)
        }

# 大文件分块上传
def upload_big_file_to_oss():
    try:
        # 获取aliyun oss client
        client = get_alicloud_oss_client()

    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }

# 下载普通文件
@retry(stop=stop_after_attempt(config["download"]["retry_attempts"]), wait=wait_fixed(config["download"]["retry_interval"]))
async def download_file(session, url, save_path):
    # 不创建目录，直接下载到指定路径
    os.makedirs(os.path.dirname(save_path), exist_ok=True)
    # 检查文件是否已经存在
    if os.path.exists(save_path):
        logger.info(f"文件已存在，跳过下载：{save_path}")
        return save_path

    # 执行下载
    async with session.get(url, timeout=aiohttp.ClientTimeout(total=config['download']['time_out'])) as response:
        if response.status != 200:
            raise Exception(f"下载失败，状态码：{response.status}, URL: {url}")
        with open(save_path, "wb") as f:
            while True:
                chunk = await response.content.read(config["download"]["chunk_size"])
                if not chunk:
                    break
                f.write(chunk)
    logger.info(f"文件下载成功：{save_path}")
    return save_path

# 下载 m3u8 文件并合并为 MP4
@retry(stop=stop_after_attempt(config["download"]["retry_attempts"]), wait=wait_fixed(config["download"]["retry_interval"]))
async def download_m3u8(session, url, save_path):
    temp_dir = config["m3u8"]["temp_path"]
    os.makedirs(temp_dir, exist_ok=True)
    temp_file = os.path.join(temp_dir, f"{Path(save_path).stem}.m3u8")
    # 下载 m3u8 文件
    async with session.get(url, timeout=aiohttp.ClientTimeout(total=config['download']['time_out'])) as response:
        if response.status != 200:
            raise Exception(f"m3u8 文件下载失败，状态码：{response.status}, URL: {url}")
        with open(temp_file, "wb") as f:
            f.write(await response.read())
    # 使用 ffmpeg 合并流媒体为 MP4
    command = f"{config['m3u8']['ffmpeg_path']} -i {url} -c copy {save_path}"
    subprocess.run(command, shell=True, check=True)
    logger.info(f"m3u8 文件下载并合并成功：{save_path}")
    return save_path

# 上传文件到 OSS
def upload_to_oss(local_path, oss_path):
    if not config["oss"]["enable_upload"]:
        logger.info(f"OSS 上传已禁用，跳过上传：{oss_path}")
        return
    result = upload_simple_file_to_oss(local_path, oss_path)
    logger.info(f"文件上传成功：{oss_path}")
    return result


# 更新数据库记录图片宽和高
def update_resource_width_height(resource_id, width, height):
    if not config["excel"]["enable_update_db"]:
        logger.info(f"update DB width height 已禁用，跳过：{resource_id}")
        return

    conn = get_db_connection()
    try:
        with conn.cursor() as cursor:
            sql = "UPDATE dim_resource_info SET width=%s, height=%s WHERE resource_id=%s"
            cursor.execute(sql, (width, height, resource_id))
            conn.commit()
    except Exception as e:
        logger.error(f"更新数据库文档宽和高失败：{e}")
    finally:
        conn.close()

# 更新 Excel 进度
def update_excel_progress(data):
    progress_file = config["excel"]["progress_file"]
    if not os.path.exists(progress_file):
        wb = Workbook()
        ws = wb.active
        ws.append(["ID", "资源ID", "资源名称", "文件名", "资源路径", "资源类型", "创建时间", "状态", "OSS完整路径", "备注"])
        wb.save(progress_file)
    wb = load_workbook(progress_file)
    ws = wb.active
    for row in data:
        ws.append(row)
    wb.save(progress_file)
    logger.info("Excel 进度已更新")

# 处理 resource_path，去掉域名并保持目录结构
def process_resource_path(resource_path):
    parsed_url = urlparse(resource_path)
    path = parsed_url.path.lstrip('/')  # 去掉开头的斜杠
    return path

# 异步任务处理
async def process_material(material, session, progress_bar):
    try:
        resource_path = material["resource_path"]
        resource_name = material["resource_name"]
        resource_type = material["resource_type_tag"]
        resource_id = material["resource_id"]
        create_time = material["create_time"]

        # 校验 resource_path 的合法性
        if not is_valid_resourcepath(resource_path):
            reason = "资源路径非法"
            logger.error(f"{reason}：{resource_path}")
            update_excel_progress([[material["id"], resource_id, resource_name, "", resource_path, resource_type, create_time, "失败", "", reason]])
            update_resource_status(resource_id=resource_id, is_downloaded=0, error_message=reason)
            return

        file_suff = os.path.splitext(resource_name)[1].lower()

        # 处理 resource_path，去掉域名并保持目录结构
        oss_path = process_resource_path(resource_path)
        file_name = os.path.basename(oss_path)

        # 获取文件的后缀名
        _, ext = os.path.splitext(file_name)
        # 检查后缀名是否为空
        if ext == '':
            full_file_name = f"{file_name}{file_suff}"
        else:
            if file_name.endswith('.m3u8'):
                if file_suff == '':
                    full_file_name = file_name[:-5] + ".mp4"
                else:
                    if file_suff in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.m4v', '.vob', '.flv', '.rmvb', '.ts', '.dv', '.f4v', '.3gp', '.webm']:
                        full_file_name = file_name[:-5] + file_suff
                    elif file_suff in ['.mp3', '.wav', '.aac', '.flac', '.ogg', '.wma', '.m4a', '.aiff', '.aif', '.m4a', '.opus', '.ac3', '.dts', '.ra']:
                        full_file_name = file_name[:-5] + file_suff
                    else:
                        full_file_name = file_name[:-5] + ".mp4"
            else:
                if file_suff == '':
                    full_file_name = file_name + ".mp4"
                else:
                    full_file_name = file_name

        # 本地下载文件路径（保持文件名不变）
        save_path = os.path.join(config["download"]["save_path"], full_file_name)

        # 根据 resource_path 后缀判断是否为 m3u8 文件
        file_extension = os.path.splitext(resource_path)[1].lower()
        if file_extension == ".m3u8":
            await download_m3u8(session, resource_path, save_path)
        elif resource_type == "视频":
            await download_file(session, resource_path, save_path)
        elif resource_type == "文档":
            await download_file(session, resource_path, save_path)
        elif resource_type == "图片":
            await download_file(session, resource_path, save_path)
        else:
            logger.info(f"不需要下载的第三方连接:{resource_path}")

        download_time = datetime.now()
        # 确认下载完整后更新数据库
        #update_resource_status(resource_id=resource_id, is_downloaded=1, download_time=datetime.now())

        # 上传到 OSS
        if resource_type == "视频":
            if file_extension == '.m3u8':
                # 对于视频，保持原有目录结构
                upload_path = oss_path[:-5] + ".mp4"
            else:
                # 对于视频，保持原有目录结构
                upload_path = oss_path
        elif resource_type == "图片":
            if file_extension == '':
                # 对于附件，直接使用文件名
                upload_path = os.path.basename(oss_path) + file_suff
            else:
                # 对于附件，直接使用文件名
                upload_path = oss_path
        elif resource_type == "文档":
            # 对于附件，直接使用文件名
            upload_path = os.path.basename(oss_path) + file_suff
        else:
            # 对于附件，直接使用文件名
            upload_path = os.path.basename(oss_path) + file_suff

        result = upload_to_oss(save_path, upload_path)

        # 生成 OSS 完整路径
        oss_complete_url = f"{config['oss']['base_url']}/{upload_path}"
        logger.info(f"oss_complete_url: {oss_complete_url}")

        # 更新 Excel
        update_excel_progress([[material["id"], resource_id, resource_name, file_name, resource_path, resource_type, create_time, "成功", oss_complete_url, ""]])
        # 确认上传完整后更新数据库
        update_resource_status(resource_id=resource_id, is_downloaded=1, download_time=download_time, is_upload_oss=1, upload_time=datetime.now(), oss_path=oss_complete_url)

        if resource_type == "文档" or resource_type == "图片":
            # 获取素材库文档类型资料的页数
            page_count = get_page_count(save_path, resource_name)
            logger.info(f"文档的资源id：{resource_id}，文档名称：{resource_name}，文档存储名：{save_path}，文档的页数：{page_count}")

            # 更新文档的页数
            update_resource_page_count(resource_id=resource_id, page_count=page_count)

        if resource_type == "图片":
            # 获取图片文档类型的宽度和高度
            width, height = get_image_size(save_path)

            # 更新图片的高和宽
            update_resource_width_height(resource_id=resource_id, width=width, height=height)

    except Exception as e:
        logger.error(f"处理失败：{material}, 错误信息：{e}")
        update_excel_progress([[material["id"], material["resource_id"], material["resource_name"], "", material["resource_path"], material["resource_type_tag"], material["create_time"], "失败", "", str(e)]])
        update_resource_status(resource_id=material["resource_id"], is_downloaded=0, error_message=str(e))
    finally:
        progress_bar.update(1)  # 更新进度条

# 主函数
async def main():
    materials = fetch_material_data()
    total = len(materials)
    semaphore = asyncio.Semaphore(config["download"]["concurrent_limit"])
    start_time = time.time()

    # 使用 tqdm 显示进度条
    with tqdm(total=total, desc="下载进度", unit="文件") as progress_bar:
        async with aiohttp.ClientSession() as session:
            tasks = []
            for material in materials:
                tasks.append(process_material_with_semaphore(material, session, semaphore, progress_bar))
            await asyncio.gather(*tasks)

    end_time = time.time()
    elapsed_time = end_time - start_time
    logger.info(f"所有任务完成，总耗时：{elapsed_time:.2f} 秒")

# 限制并发数
async def process_material_with_semaphore(material, session, semaphore, progress_bar):
    async with semaphore:
        await process_material(material, session, progress_bar)

# 获取pdf页数
def get_pdf_page_count(pdf_path):
    """获取PDF文件页数"""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            return len(reader.pages)
    except Exception as e:
        logger.error(f"读取PDF文件失败: {e}")
        return None

# 获取word页数
def get_word_page_count(word_path):
    """获取Word文件页数"""
    try:
        doc = Document(word_path)
        # Word文档的页数通常不能直接获取，这里计算段落数作为替代
        return len(doc.paragraphs)  # 注意：这并不是准确的页数
    except Exception as e:
        logger.error(f"读取Word文件失败: {e}")
        return None

# 获取ppt页数
def get_ppt_page_count(ppt_path):
    """获取PPT文件页数"""
    try:
        presentation = Presentation(ppt_path)
        return len(presentation.slides)
    except Exception as e:
        logger.error(f"读取PPT文件失败: {e}")
        return None

# 获取excel页数 默认：1
def get_excel_page_count(excel_path):
    """获取Excel文件页数"""
    # 假设Excel文件的页数为1
    return 1

# 获取图片页数 默认：1
def get_image_page_count(image_path):
    """获取图片文件页数"""
    # 假设图片文件的页数为1
    return 1

# 获取wps文档页数
def get_wps_page_count(wps_path):
    """获取WPS文档页数"""
    try:
        doc = Document(wps_path)  # 可能需要安装wps-office或相关库
        return len(doc.paragraphs)  # 注意：这并不是准确的页数
    except Exception as e:
        logger.error(f"读取WPS文件失败: {e}")
        return None

# 获取et表格页数   默认：1
def get_et_page_count(et_path):
    """获取WPS表格页数"""
    return 1  # 假设WPS表格的页数为1

# 获取dps演示文档页数
def get_dps_page_count(dps_path):
    """获取WPS演示文稿页数"""
    try:
        presentation = Presentation(dps_path)  # 可能需要安装wps-office或相关库
        return len(presentation.slides)
    except Exception as e:
        logger.error(f"读取DPS文件失败: {e}")
        return None

# 获取文档的页数
def get_page_count(file_path, resource_name):
    """根据文件后缀判断文件类型并获取页数"""
    _, ext = os.path.splitext(resource_name.lower())

    if ext == '.pdf':
        return get_pdf_page_count(file_path)
    elif ext in ['.doc', '.docx']:
        return get_word_page_count(file_path)
    elif ext in ['.ppt', '.pptx']:
        return get_ppt_page_count(file_path)
    elif ext in ['.xls', '.xlsx']:
        return get_excel_page_count(file_path)
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.heic', '.heif', '.ico', '.tiff', '.tif', '.webp', '.svg', '.raw']:
        return get_image_page_count(file_path)
    elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.m4v', '.vob', '.flv', '.rmvb', '.ts', '.dv', '.f4v', '.3gp', '.webm']:
        return 0
    elif ext in ['.mp3', '.wav', '.aac', '.flac', '.ogg', '.wma', '.m4a', '.aiff', '.aif', '.m4a', '.opus', '.ac3', '.dts', '.ra']:
        return 0
    elif ext in ['.wps']:
        return get_wps_page_count(file_path)
    elif ext in ['.dps']:
        return get_dps_page_count(file_path)
    elif ext in ['.et']:
        return get_et_page_count(file_path)
    else:
        logger.warning(f"不支持的文件类型: {ext}")
        return 0

# 获取图片文档类型的宽度和高度
def get_image_size(image_path):
    """
    获取图片的宽度和高度。

    :param image_path: 图片文件的路径
    :return: 一个元组 (宽度, 高度)
    """
    with Image.open(image_path) as img:
        return img.size  # 返回 (宽度, 高度)

if __name__ == "__main__":
    asyncio.run(main())
